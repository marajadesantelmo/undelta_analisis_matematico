"""
Script para generar la presentación de Análisis Matemático I - Clase 1.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de colores ──────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)   # azul marino muy oscuro
ACCENT     = RGBColor(0x1F, 0x8E, 0xF0)   # azul brillante
ACCENT2    = RGBColor(0xF0, 0xA5, 0x00)   # naranja dorado
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE8)
DARK_GRAY  = RGBColor(0x1A, 0x2A, 0x40)


def set_slide_bg(slide, color: RGBColor):
    """Pone un fondo sólido a la diapositiva."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txb(slide, text, left, top, width, height,
        font_size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Agrega un TextBox simple."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, color=WHITE, bullet_color=None,
                    heading=None, heading_size=18):
    """Agrega un cuadro de texto con lista de viñetas."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    if heading:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = heading
        run.font.size = Pt(heading_size)
        run.font.bold = True
        run.font.color.rgb = ACCENT2

    for i, item in enumerate(items):
        if heading or i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        p.space_before = Pt(4)

    return tb


# ── Crear presentación 16:9 ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # layout en blanco

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Portada
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

# Barra decorativa izquierda
add_rect(slide, 0, 0, 0.35, 7.5, ACCENT)

# Barra inferior
add_rect(slide, 0, 6.8, 13.33, 0.7, DARK_GRAY)

# Título principal
txb(slide, "Análisis Matemático I",
    0.6, 1.8, 12.0, 1.4, font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Línea divisoria decorativa
add_rect(slide, 0.6, 3.3, 6.0, 0.07, ACCENT)

# Subtítulo
txb(slide, "Clase 1: Introducción y el Concepto de Función",
    0.6, 3.5, 12.0, 0.9, font_size=26, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

# Datos inferiores
txb(slide, "Ciencias de Datos – Comisión 4",
    0.6, 6.85, 9.0, 0.5, font_size=14, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
txb(slide, "Docentes: Facundo Lastra  ·  Pablo Verón",
    0.6, 5.6, 9.0, 0.55, font_size=17, bold=True, color=ACCENT2, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Hoja de Ruta
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT)

txb(slide, "¿Qué haremos hoy?",
    0.5, 0.15, 12.0, 1.0, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

items = [
    "Presentación, modalidad y condiciones de cursada.",
    "El corazón de la materia: ¿Qué es una función?",
    "El rol vital de las funciones en la Ciencia de Datos.",
    "Ejemplo práctico: Análisis de una señal.",
    "Herramientas: Uso de IA y simuladores.",
    "Reflexión: La matemática como búsqueda de patrones.",
    "Actividad interactiva de cierre.",
]
add_bullet_list(slide, items, 0.8, 1.6, 11.5, 5.5, font_size=18, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Modalidad de Cursada
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT)

txb(slide, "¿Cómo vamos a trabajar?",
    0.5, 0.15, 12.0, 1.0, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txb(slide, "(Aquí compartís pantalla con el programa de la materia)",
    0.8, 1.55, 11.5, 0.55, font_size=15, italic=True, color=ACCENT2, align=PP_ALIGN.LEFT)

items = [
    "Días y horarios de clases (teóricas y prácticas).",
    "Condiciones de regularidad y promoción.",
    "Fechas importantes (parciales).",
    "Canales de comunicación (Campus, foros, etc.).",
]
add_bullet_list(slide, items, 0.8, 2.3, 11.5, 4.5, font_size=20, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – El Concepto de Función
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT)

txb(slide, "Nuestro primer gran objetivo: La Función",
    0.5, 0.15, 12.0, 1.0, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txb(slide,
    "En matemática, una función es mucho más que una fórmula; es una regla de asignación.",
    0.8, 1.6, 11.5, 0.9, font_size=20, color=LIGHT_GRAY)

# Caja de definición destacada
add_rect(slide, 0.6, 2.7, 12.0, 2.5, DARK_GRAY)
add_rect(slide, 0.6, 2.7, 0.12, 2.5, ACCENT)

txb(slide, "Definición Clave",
    0.9, 2.8, 11.0, 0.5, font_size=18, bold=True, color=ACCENT2)
txb(slide,
    "Es una regla que a cada elemento de un conjunto de partida (llamado Dominio) "
    "le asigna un ÚNICO elemento de un conjunto de llegada (llamado Codominio o Imagen).",
    0.9, 3.35, 11.2, 1.7, font_size=19, color=WHITE)

# Diagrama simplificado con flechas de texto
txb(slide, "Dominio   →   f( )   →   Codominio / Imagen",
    1.5, 5.5, 10.0, 0.7, font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Funciones en Ciencias de Datos
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT)

txb(slide, "¿Por qué estudiamos esto en Ciencias de Datos?",
    0.5, 0.15, 12.0, 1.0, font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txb(slide,
    "El enfoque cambia radicalmente según lo que estemos haciendo. "
    "Es clave entender esta diferencia:",
    0.8, 1.55, 11.5, 0.7, font_size=18, color=LIGHT_GRAY)

# Caja izquierda – Análisis Matemático
add_rect(slide, 0.5, 2.5, 5.8, 4.2, DARK_GRAY)
add_rect(slide, 0.5, 2.5, 5.8, 0.12, ACCENT)
txb(slide, "En Análisis Matemático I",
    0.7, 2.65, 5.3, 0.6, font_size=17, bold=True, color=ACCENT)
txb(slide,
    "Ya conocemos la función y nos dedicamos a estudiar sus propiedades: "
    "sus límites, cómo crece, dónde tiene máximos o mínimos.",
    0.7, 3.35, 5.3, 3.1, font_size=17, color=WHITE)

# Caja derecha – Ciencia de Datos
add_rect(slide, 7.0, 2.5, 5.8, 4.2, DARK_GRAY)
add_rect(slide, 7.0, 2.5, 5.8, 0.12, ACCENT2)
txb(slide, "En Ciencia de Datos",
    7.2, 2.65, 5.3, 0.6, font_size=17, bold=True, color=ACCENT2)
txb(slide,
    "Tenemos un montón de datos (ejemplos del mundo real) y nuestro objetivo "
    "es intentar aproximar funciones desconocidas que expliquen esos datos.",
    7.2, 3.35, 5.3, 3.1, font_size=17, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Análisis vs. Machine Learning
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT)

txb(slide, "Cambio de Paradigma",
    0.5, 0.15, 12.0, 1.0, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Bloque Análisis
add_rect(slide, 0.5, 1.55, 5.8, 2.8, DARK_GRAY)
add_rect(slide, 0.5, 1.55, 5.8, 0.1, ACCENT)
txb(slide, "Análisis Matemático",
    0.7, 1.7, 5.3, 0.55, font_size=18, bold=True, color=ACCENT)
txb(slide, "Tenemos la regla y los datos de entrada,\ncalculamos la salida.",
    0.7, 2.35, 5.3, 0.85, font_size=16, color=LIGHT_GRAY)
txb(slide, "y = f(x)",
    1.2, 3.1, 4.5, 0.8, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bloque ML
add_rect(slide, 7.0, 1.55, 5.8, 2.8, DARK_GRAY)
add_rect(slide, 7.0, 1.55, 5.8, 0.1, ACCENT2)
txb(slide, "Machine Learning",
    7.2, 1.7, 5.3, 0.55, font_size=18, bold=True, color=ACCENT2)
txb(slide,
    "Tenemos los datos de entrada y salida,\ny queremos que la computadora \"aprenda\" la regla.",
    7.2, 2.35, 5.3, 0.85, font_size=16, color=LIGHT_GRAY)
txb(slide, "(x, y)  →  f̂",
    7.5, 3.1, 4.5, 0.8, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Conclusión
add_rect(slide, 0.5, 5.0, 12.3, 1.9, RGBColor(0x1F, 0x3A, 0x5F))
add_rect(slide, 0.5, 5.0, 12.3, 0.1, ACCENT2)
txb(slide,
    "La Ciencia de Datos es, en el fondo:\nTeoría de Aproximación + Optimización.",
    1.0, 5.15, 11.3, 1.6, font_size=24, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Ejemplo Práctico
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT)

txb(slide, "Analizando una Señal de Radio",
    0.5, 0.15, 12.0, 1.0, font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Placeholder para imagen
add_rect(slide, 0.5, 1.55, 7.5, 5.6, RGBColor(0x12, 0x25, 0x3A))
txb(slide, "[ Insertar aquí el gráfico\nAmplitud en función del tiempo ]",
    0.5, 3.0, 7.5, 2.0, font_size=15, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Preguntas
questions = [
    "a.  ¿Cuál es la amplitud máxima alcanzada?",
    "b.  ¿En qué instantes la señal tuvo 3 000 cm de amplitud?",
    "c.  ¿Qué amplitud se produce a los 2 segundos?",
    "d.  ¿En qué intervalos la amplitud se mantuvo constante?",
]
txb(slide, "Preguntas para debatir:", 8.3, 1.6, 4.8, 0.5, font_size=16, bold=True, color=ACCENT2)
add_bullet_list(slide, questions, 8.3, 2.2, 4.8, 4.8, font_size=15, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – IA como herramienta de estudio
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT)

txb(slide, "Usando IA para visualizar matemática",
    0.5, 0.15, 12.0, 1.0, font_size=33, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txb(slide,
    "La Inteligencia Artificial no es solo para escribir textos, "
    "es una excelente herramienta para explorar conceptos matemáticos.",
    0.8, 1.55, 11.5, 0.9, font_size=19, color=LIGHT_GRAY)

add_rect(slide, 0.5, 2.7, 12.3, 3.8, DARK_GRAY)
add_rect(slide, 0.5, 2.7, 0.12, 3.8, ACCENT2)
txb(slide, "Actividad en vivo",
    0.85, 2.85, 11.5, 0.55, font_size=20, bold=True, color=ACCENT2)
txb(slide,
    "En este momento puedes abrir ChatGPT o Gemini, compartir pantalla y "
    "mostrarles cómo pedirle a la IA que:\n\n"
    "  •  Genere ejemplos reales de funciones lineales.\n\n"
    "  •  Explique cómo se aplica una función lineal en un modelo predictivo sencillo.",
    0.85, 3.5, 11.5, 2.8, font_size=18, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Cambiando la Perspectiva
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT)

txb(slide, "La matemática como lente para ver el mundo",
    0.5, 0.15, 12.0, 1.0, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

items = [
    "La matemática se trata de buscar patrones, representarlos, tomar supuestos y ver sus efectos.",
    "Es un lenguaje para representar la realidad (como el braille o una partitura musical).",
    "Conocer implica cambiar la perspectiva: una ecuación nos permite ver un mismo problema desde diferentes ángulos.",
    "Entender es poder abstraer.\n     Ejemplo: Reconocer el patrón de una letra "R", sin importar su tipografía, tamaño o color.\n     La matemática nos da el poder para hacer eso con los datos.",
]
add_bullet_list(slide, items, 0.8, 1.6, 11.5, 5.5, font_size=17, color=LIGHT_GRAY)

txb(slide, "(Basados en la charla de Roger Antonsen)",
    8.5, 6.9, 4.5, 0.4, font_size=11, italic=True, color=ACCENT, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Actividad de Cierre
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 1.3, DARK_GRAY)
add_rect(slide, 0, 1.28, 13.33, 0.06, ACCENT2)

txb(slide, "¡Manos a la obra!",
    0.5, 0.15, 12.0, 1.0, font_size=40, bold=True, color=ACCENT2, align=PP_ALIGN.LEFT)

txb(slide, "Vamos a poner en práctica lo que charlamos hoy.",
    0.8, 1.55, 11.5, 0.65, font_size=20, color=LIGHT_GRAY)

txb(slide, "Instrucciones:", 0.8, 2.45, 11.5, 0.5, font_size=20, bold=True, color=WHITE)
txb(slide,
    "1.  Ingresen desde sus computadoras o celulares a la aplicación interactiva.\n\n"
    "2.  Exploren los parámetros y vean cómo cambian los gráficos.",
    0.8, 3.05, 11.5, 1.6, font_size=19, color=LIGHT_GRAY)

# Caja con enlace
add_rect(slide, 0.5, 4.9, 12.3, 1.5, RGBColor(0x08, 0x3A, 0x6A))
add_rect(slide, 0.5, 4.9, 12.3, 0.1, ACCENT2)
txb(slide, "https://undelta-ami.streamlit.app/",
    0.5, 5.1, 12.3, 1.0, font_size=26, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# Barra inferior
add_rect(slide, 0, 6.8, 13.33, 0.7, DARK_GRAY)
txb(slide, "Análisis Matemático I  ·  Ciencias de Datos – Comisión 4",
    0.5, 6.85, 12.0, 0.4, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Guardar ────────────────────────────────────────────────────────────────────
output_path = "Clase1_Analisis_Matematico_I.pptx"
prs.save(output_path)
print(f"Presentación guardada en: {output_path}")
